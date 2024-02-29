
from pptx import Presentation
from pptx.util import Cm, Pt

from lib.google_images import GoogleImage
from lib.utils import *
from lib.config import *


def set_font(element, name: str = None, size: int = None, bold: bool = None, italic: bool = None):
    t = type(element).__name__
    if t == "SlidePlaceholder":
        element = element.text_frame.paragraphs[0]
    font = element.font
    if name is not None:
        font.name = name
    if size is not None:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if italic is not None:
        font.italic = italic


def make_pptx(pptx_template_path,
              topic: str,
              speaker: str,
              contents: list,
              title_font={},
              speaker_font={},
              slide_title_font={},
              content_font={}):
    
    if pptx_template_path:
        with open(pptx_template_path, "rb") as f:
            prs = Presentation(f)
    else:
        prs = Presentation()
    prs.core_properties.language = "de"
    
    num_content_slides = len(contents)
    title_slide_layout = random.choice([lay for lay in prs.slide_layouts if lay.name.endswith("title")])
    bullet_slide_layouts = sample_minimal_repitions([lay for lay in prs.slide_layouts if lay.name.endswith("bullets")], num_content_slides)
    bullet_img_slide_layouts = sample_minimal_repitions([lay for lay in prs.slide_layouts if lay.name.endswith("bullets_img")], num_content_slides)
    
    slide = prs.slides.add_slide(title_slide_layout)
    e_title = slide.shapes.title
    e_title.text = topic
    set_font(e_title, **title_font)
    e_speaker = slide.placeholders[1]

    e_speaker.text = speaker
    set_font(e_speaker, **speaker_font)
    
    for i, content in enumerate(contents):
        img: GoogleImage = content.get("img", None)
        
        if img:
            layout = bullet_img_slide_layouts[i]
        else:
            layout = bullet_slide_layouts[i]
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        
        e_slide_title = shapes.title
        e_slide_title.text = content["title"]
        set_font(e_slide_title.text_frame.paragraphs[0], **slide_title_font)

        if img is not None:
            # print([s.name for s in slide.shapes])
            # container_shape = [s for s in slide.shapes if s.name == "img_group"][0]
            ext = str(img.local_path).split(".")[-1]
            if ext.lower() not in ["bmp", "gif", "jpg", "jpeg", "png"]:
                continue
            max_width, max_height = 11.5, 12.57
            x, y = 20.3, 5
            aspect = img.width / img.height
            if aspect < max_width / max_height:
                size = {"height": Cm(max_height)}
                w = aspect * max_height
                # align right
                x += max_width - w
            else:
                size = {"width": Cm(max_width)}
                h = max_width / aspect
                # align vertical middle
                y += (max_height - h) / 2
            try:
                slide.shapes.add_picture(str(img.local_path), left=Cm(x), top=Cm(y), **size)
                # container_shape.add_picture(str(img.local_path), left=Cm(0), top=Cm(0), **size)
            except Exception:
                pass
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        for j, (pstyle, text, *args) in enumerate(content["content"]):
            if j > 0:
                p = tf.add_paragraph()
            if pstyle == "bullet" and len(args):
                p.level = args[0]
            set_font(p, **content_font)

            parts = [("normal" if k % 2 == 0 else "bold", s) for k, s in enumerate(text.split("**"))]

            for rstyle, s in parts:
                run = p.add_run()
                run.text = s
                if rstyle == "bold":
                    set_font(run, bold=True)
    
    output_path = PPTX_DIR / f"pptx-{NOW()}-{ESCAPE_PATH(speaker)}-{ESCAPE_PATH(topic)}.pptx"
    # output_path = PPTX_DIR / f"pptx-{NOW()}-{ESCAPE_PATH(speaker)}.pptx"
    prs.save(output_path)
    return output_path



if __name__ == "__main__":
    pass
    # Test make_pptx
    # speaker = "Johann Wolfram von .Brüste"
    # topic = "Die Otten"
    # _topic, contents = parse_md_outline(open(TEMPLATE_DIR / "markdown" / "Alternative Energien.md", "r", encoding="utf8").read())

    # imgs = [
    #     GoogleImage(result_index=1, url='https://i.ytimg.com/vi/ZVJEwmazSbM/maxresdefault.jpg', context_url='https://m.youtube.com/watch?v=ZVJEwmazSbM', accessed=datetime(2024, 2, 29, 1, 55, 23, 55177), request={'title': 'Google Custom Search - Fischotter', 'totalResults': '13300000', 'searchTerms': 'Fischotter', 'count': 3, 'startIndex': 1, 'inputEncoding': 'utf8', 'outputEncoding': 'utf8', 'safe': 'active', 'cx': '12938424aad76437a', 'gl': 'de', 'hl': 'de', 'searchType': 'image'}, search_parameters={'key': 'AIzaSyBbbw_-i5Hi7XisOMdGMbO-Gsy1pSWouUE', 'cx': '12938424aad76437a', 'searchType': 'image', 'q': 'Fischotter', 'num': 3, 'gl': 'de', 'hl': 'de', 'safe': 'active'}, width=1280, height=720, byte_size=132030, file_format='image/jpeg', title=None, query='Fischotter', ext='.jpg', downloaded=True, local_path=Path('D:/GitProjects/gpt-karaok/gpt-karaok-backend/tmp/img/google-img-20240229-015523-002-17e4-Fischotter.jpg'))
    # ]
    # for i, slide in enumerate(contents):
    #     if i < len(imgs):
    #         slide["img"] = imgs[i]
    
    # pptx_path = make_pptx(pptx_template_path=PPTX_TEMPLATE_DIR / "template-white-16-9-colorful-wild-animations.pptx",
    #                       contents=contents, speaker=speaker, topic=topic)
    # show_presentation_blocking(pptx_path)